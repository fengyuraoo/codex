from __future__ import annotations

import hashlib
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from . import database
from .classifier import classify, infer_tags, weighted_match, MATERIAL_RULES, PROJECT_RULES
from .history_service import finish_scan_batch, start_scan_batch
from .import_service import file_hash, image_preview
from .models import MaterialRecord
from .paths import CONTENT_DIR, DATA_FOLDERS, ROOT
from .utils import atomic_write_text


SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".jpg", ".jpeg", ".png", ".docx", ".pptx"}


def infer_source_mode(path: Path) -> str:
    rel = str(path.relative_to(ROOT)).replace("\\", "/")
    if rel.startswith("data/examples/"):
        return "demo"
    if rel.startswith("data/inbox/"):
        return "user"
    if rel.startswith("data/library/"):
        return "imported"
    return "unknown"


def stable_id(path: Path) -> str:
    rel_path = str(path.relative_to(ROOT)).replace("\\", "/").lower()
    digest = hashlib.sha1(rel_path.encode("utf-8")).hexdigest()[:12]
    return f"mat-{digest}"


def count_words(text: str) -> int:
    if not text.strip():
        return 0
    chinese_chars = re.findall(r"[\u4e00-\u9fff]", text)
    latin_words = re.findall(r"[A-Za-z0-9_]+", text)
    return len(chinese_chars) + len(latin_words)


def read_text(path: Path) -> tuple[str, str, str]:
    errors = []
    for encoding in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return path.read_text(encoding=encoding), "parsed", ""
        except UnicodeDecodeError as exc:
            errors.append(f"{encoding}: {exc}")
    return "", "text_unavailable", "; ".join(errors)


def extract_pdf_text(path: Path) -> tuple[str, str, str]:
    try:
        import fitz  # type: ignore

        with fitz.open(path) as doc:
            return "\n".join(page.get_text() for page in doc).strip(), "parsed", ""
    except ImportError:
        pass
    except Exception as exc:
        return "", "parse_failed", f"PyMuPDF error: {exc}"
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip(), "parsed", ""
    except ImportError:
        return "", "pdf_text_unavailable", "Install PyMuPDF or pypdf to parse PDF text."
    except Exception as exc:
        return "", "parse_failed", f"pypdf error: {exc}"


def extract_docx_text(path: Path) -> tuple[str, str, str]:
    try:
        import docx  # type: ignore

        document = docx.Document(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip(), "parsed", ""
    except ImportError:
        return "", "docx_text_unavailable", "Install python-docx to parse DOCX text."
    except Exception as exc:
        return "", "parse_failed", f"python-docx error: {exc}"


def extract_pptx_text(path: Path) -> tuple[str, str, str]:
    try:
        from pptx import Presentation  # type: ignore

        presentation = Presentation(str(path))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts).strip(), "parsed", ""
    except ImportError:
        return "", "pptx_text_unavailable", "Install python-pptx to parse PPTX text."
    except Exception as exc:
        return "", "parse_failed", f"python-pptx error: {exc}"


def image_metadata(path: Path) -> tuple[str, str, str, int, int]:
    try:
        from PIL import Image  # type: ignore

        with Image.open(path) as image:
            return f"Image size: {image.width}x{image.height}. 待接入图像理解模型。", "metadata_only", "", image.width, image.height
    except ImportError:
        return "Image metadata unavailable. 待接入图像理解模型。", "image_metadata_unavailable", "Install Pillow to read image dimensions.", 0, 0
    except Exception as exc:
        return "Image parse failed. 待接入图像理解模型。", "parse_failed", f"Pillow error: {exc}", 0, 0


def parse_file_payload(path: Path) -> tuple[str, str, str, int, int]:
    extension = path.suffix.lower()
    if extension in {".md", ".txt"}:
        text, status, error = read_text(path)
        return text, status, error, 0, 0
    if extension == ".pdf":
        text, status, error = extract_pdf_text(path)
        return text, status, error, 0, 0
    if extension == ".docx":
        text, status, error = extract_docx_text(path)
        return text, status, error, 0, 0
    if extension == ".pptx":
        text, status, error = extract_pptx_text(path)
        return text, status, error, 0, 0
    if extension in {".jpg", ".jpeg", ".png"}:
        return image_metadata(path)
    return "", "unsupported", "Unsupported file type.", 0, 0


def parse_file_text(path: Path) -> tuple[str, str, str]:
    text, status, error, _, _ = parse_file_payload(path)
    return text, status, error


def parse_material(path: Path, source_folder: Path, scan_batch_id: str = "") -> MaterialRecord:
    stat = path.stat()
    full_text, parse_status, parse_error, image_width, image_height = parse_file_payload(path)
    material_id = stable_id(path)
    digest = file_hash(path)
    content_full_path = ""
    if full_text.strip() and parse_status == "parsed":
        content_path = CONTENT_DIR / f"{material_id}.txt"
        atomic_write_text(content_path, full_text)
        content_full_path = str(content_path.relative_to(ROOT)).replace("\\", "/")

    combined = f"{path.name} {full_text[:1000]}"
    material_type, _, _ = weighted_match(combined, MATERIAL_RULES)
    project_guess, _, _ = weighted_match(combined, PROJECT_RULES)
    record = MaterialRecord(
        id=material_id,
        filename=path.name,
        path=str(path.relative_to(ROOT)).replace("\\", "/"),
        extension=path.suffix.lower().lstrip("."),
        size=stat.st_size,
        modified_time=datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
        source_folder=str(source_folder.relative_to(ROOT)).replace("\\", "/"),
        content_preview=full_text[:1000].strip(),
        content_full_path=content_full_path,
        parse_status=parse_status,
        parse_error=parse_error,
        word_count=count_words(full_text if parse_status == "parsed" else ""),
        material_type=material_type,
        project_guess=project_guess,
        tags=infer_tags(combined),
        file_hash=digest,
        scan_batch_id=scan_batch_id,
        is_duplicate=database.hash_seen_elsewhere(material_id, digest),
        image_preview_path=image_preview(path, material_id) if path.suffix.lower() in {".jpg", ".jpeg", ".png"} else "",
        image_width=image_width,
        image_height=image_height,
        image_note="当前未接入图像理解，可手动填写图片说明。" if path.suffix.lower() in {".jpg", ".jpeg", ".png"} else "",
        source_mode=infer_source_mode(path),
    )
    return classify(record, full_text)


def scan_library() -> list[MaterialRecord]:
    CONTENT_DIR.mkdir(parents=True, exist_ok=True)
    records = []
    batch_id, _ = start_scan_batch()
    stats = {"total_files": 0, "new_files": 0, "updated_files": 0, "duplicate_files": 0, "failed_files": 0}
    for folder in DATA_FOLDERS:
        folder.mkdir(parents=True, exist_ok=True)
        for path in sorted(folder.rglob("*")):
            if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
                stats["total_files"] += 1
                material_id = stable_id(path)
                existed = database.material_exists(material_id)
                record = parse_material(path, folder, batch_id)
                if record.is_duplicate:
                    stats["duplicate_files"] += 1
                elif existed:
                    stats["updated_files"] += 1
                else:
                    stats["new_files"] += 1
                if record.parse_status == "parse_failed":
                    stats["failed_files"] += 1
                records.append(record)
    finish_scan_batch(batch_id, stats)
    return records
